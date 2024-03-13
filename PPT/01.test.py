# 导包
from pptx import Presentation

# 创建空白演示文稿
prs = Presentation()
# 添加标题布局的幻灯片
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
# 设置标题和副标题
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"
# 保存
prs.save('test.pptx')

