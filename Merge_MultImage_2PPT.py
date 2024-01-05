#!/usr/bin/python
# -*- encoding: utf-8 -*-
'''
@File    :   Merge_MultImage_2PPT.py
@Time    :   2023/12/12 14:23:22
@Author  :   Liu.Bo 
@Version :   1.0.0.0
@Contact :   liubo4@genomics.cn/614347533@qq.com
@WebSite :   http://www.ben-air.cn/
'''

import logging
from logging.handlers import RotatingFileHandler
from argparse import ArgumentParser
import os
# pip install python-pptx
from pptx import Presentation

program = 'Merge_MultImage_2PPT.py'
version = '1.0.0.0'

parser = ArgumentParser(prog=program)
parser.add_argument('-indir' , dest='input_dir'   , action='store', type=str, help='input  File',default="")
parser.add_argument('-list', dest='input_dir_list', action='store', type=str, help='output File', default="")
args = parser.parse_args()

def create_ppt_from_images(directory, output_file):
    # 创建 PowerPoint 对象
    prs = Presentation()
    prs.slide_height = 9144000 #设置slide的高
    prs.slide_width = 16144000 #设置slide的宽
    # 遍历目录中的图片文件
    for filename in os.listdir(directory):
        if filename.endswith(".jpg") or filename.endswith(".png"):
            # 创建新的幻灯片
            slide_layout = prs.slide_layouts[1]  # 使用第二个默认布局（标题和内容）
            slide = prs.slides.add_slide(slide_layout)

            # 添加图片到幻灯片
            image_path = os.path.join(directory, filename)
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # 保存 PPT 文件
    prs.save(output_file)

def main():
    if len(args.input_dir) >1:
        create_ppt_from_images(args.input_dir, args.input_dir+".pptx")
    elif len(args.input_dir_list) >0:
        Dirlist=open(args.input_dir_list,"r", encoding="utf-8")
        for dir in Dirlist.readlines():
            dir=dir.strip()
            create_ppt_from_images(dir,dir+".pptx")
    else:
        print("""Usage: 
python Merge_MultImage_2PPT.py
    -indir  a Direction contain images 
    -list   a list of directions contain images
list Demo>>>
    /root/dirA
    /root/dirB
<<<

              """)

if __name__ == '__main__':
    main()


